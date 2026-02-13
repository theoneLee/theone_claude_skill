#!/usr/bin/env python3
"""
PPT 生成脚本：将 Markdown 手写稿转换为精美的 PPTX 演示文稿。

用法:
    python3 generate_ppt.py --input draft.md --output slides.pptx --title "标题" --theme tech_blue

依赖:
    pip install python-pptx Pillow
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("错误: 请先安装 python-pptx: pip install python-pptx Pillow")
    sys.exit(1)


# ── 配色方案 ──────────────────────────────────────────────────────────────────

THEMES = {
    "tech_blue": {
        "name": "科技蓝",
        "primary": RGBColor(0x1A, 0x73, 0xE8),
        "secondary": RGBColor(0x34, 0xA8, 0x53),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x20, 0x20, 0x20),
        "text_light": RGBColor(0x5F, 0x63, 0x68),
        "accent": RGBColor(0xEA, 0x43, 0x35),
    },
    "vibrant_orange": {
        "name": "活力橙",
        "primary": RGBColor(0xFF, 0x6D, 0x00),
        "secondary": RGBColor(0xFF, 0xD6, 0x00),
        "bg": RGBColor(0xFF, 0xF8, 0xE1),
        "text": RGBColor(0x33, 0x33, 0x33),
        "text_light": RGBColor(0x75, 0x75, 0x75),
        "accent": RGBColor(0xE6, 0x51, 0x00),
    },
    "fresh_green": {
        "name": "清新绿",
        "primary": RGBColor(0x00, 0xC8, 0x53),
        "secondary": RGBColor(0x00, 0xBF, 0xA5),
        "bg": RGBColor(0xE8, 0xF5, 0xE9),
        "text": RGBColor(0x1B, 0x5E, 0x20),
        "text_light": RGBColor(0x4C, 0xAF, 0x50),
        "accent": RGBColor(0x00, 0x96, 0x88),
    },
    "deep_purple": {
        "name": "深邃紫",
        "primary": RGBColor(0x62, 0x00, 0xEA),
        "secondary": RGBColor(0xAA, 0x00, 0xFF),
        "bg": RGBColor(0xF3, 0xE5, 0xF5),
        "text": RGBColor(0x31, 0x1B, 0x92),
        "text_light": RGBColor(0x7C, 0x4D, 0xFF),
        "accent": RGBColor(0xE0, 0x40, 0xFB),
    },
    "business_grey": {
        "name": "商务灰",
        "primary": RGBColor(0x37, 0x47, 0x4F),
        "secondary": RGBColor(0x60, 0x7D, 0x8B),
        "bg": RGBColor(0xEC, 0xEF, 0xF1),
        "text": RGBColor(0x26, 0x32, 0x38),
        "text_light": RGBColor(0x78, 0x90, 0x9C),
        "accent": RGBColor(0xFF, 0x6F, 0x00),
    },
}

# 幻灯片尺寸 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def parse_markdown(md_path: str) -> list[dict]:
    """解析 Markdown 文件，提取结构化内容。

    返回列表，每项包含:
        - type: 'h1' | 'h2' | 'content'
        - title: 标题文字
        - bullets: 要点列表
    """
    text = Path(md_path).read_text(encoding="utf-8")
    sections = []
    current = None

    for line in text.splitlines():
        line = line.rstrip()

        if line.startswith("# "):
            if current:
                sections.append(current)
            current = {"type": "h1", "title": line[2:].strip(), "bullets": []}

        elif line.startswith("## "):
            if current:
                sections.append(current)
            current = {"type": "h2", "title": line[3:].strip(), "bullets": []}

        elif line.startswith("### "):
            if current:
                sections.append(current)
            current = {"type": "h2", "title": line[4:].strip(), "bullets": []}

        elif re.match(r"^[-*]\s+", line):
            bullet = re.sub(r"^[-*]\s+", "", line).strip()
            if bullet and current:
                current["bullets"].append(bullet)

        elif line.strip() and current and not line.startswith("```"):
            # 普通文本行作为要点
            if current["bullets"] or current["type"] != "h1":
                current["bullets"].append(line.strip())

    if current:
        sections.append(current)

    return sections


def set_slide_bg(slide, color: RGBColor):
    """设置幻灯片背景色。"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_colored_bar(slide, theme: dict, position="top"):
    """添加装饰性色条。"""
    if position == "top":
        left, top = Emu(0), Emu(0)
        width, height = SLIDE_WIDTH, Inches(0.15)
    else:
        left, top = Emu(0), SLIDE_HEIGHT - Inches(0.08)
        width, height = SLIDE_WIDTH, Inches(0.08)

    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["primary"]
    shape.line.fill.background()


def create_cover_slide(prs, title: str, subtitle: str, theme: dict):
    """创建封面页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, theme["bg"])

    # 主色块背景（左侧 40%）
    shape = slide.shapes.add_shape(
        1, Emu(0), Emu(0), Inches(5.333), SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme["primary"]
    shape.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 副标题
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p2.space_before = Pt(20)

    return slide


def create_section_slide(prs, title: str, theme: dict):
    """创建章节过渡页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["primary"])

    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(prs, title: str, bullets: list[str], theme: dict, image_path: str = None):
    """创建内容页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])
    add_colored_bar(slide, theme, "top")
    add_colored_bar(slide, theme, "bottom")

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]

    # 内容区域宽度取决于是否有图片
    content_width = Inches(7) if image_path else Inches(11)

    # 要点
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), content_width, Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets[:6]):  # 最多 6 个要点
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"● {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = theme["text"]
        p.space_before = Pt(12)
        p.space_after = Pt(4)

    # 如果有图片，添加到右侧
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(
            image_path, Inches(8.5), Inches(1.8), Inches(4), Inches(4)
        )

    return slide


def create_ending_slide(prs, theme: dict):
    """创建结尾页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["primary"])

    txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "感谢观看"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "欢迎关注 · 点赞 · 评论"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(30)

    return slide


def generate_ppt(input_path: str, output_path: str, title: str, theme_name: str = "tech_blue"):
    """主生成函数：Markdown → PPTX。"""
    theme = THEMES.get(theme_name, THEMES["tech_blue"])
    sections = parse_markdown(input_path)

    if not sections:
        print(f"错误: 未能从 {input_path} 解析出任何内容")
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 封面
    main_title = title or (sections[0]["title"] if sections else "Vlog")
    create_cover_slide(prs, main_title, "Vlog 演示文稿", theme)

    # 内容页
    for section in sections:
        if section["type"] == "h1" and section != sections[0]:
            create_section_slide(prs, section["title"], theme)
        elif section["bullets"]:
            create_content_slide(prs, section["title"], section["bullets"], theme)
        elif section["type"] == "h2":
            create_section_slide(prs, section["title"], theme)

    # 结尾
    create_ending_slide(prs, theme)

    # 保存
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"✅ PPT 已生成: {output}")
    print(f"   共 {len(prs.slides)} 页")
    print(f"   配色方案: {theme['name']}")


def main():
    parser = argparse.ArgumentParser(description="将 Markdown 手写稿转换为 PPT")
    parser.add_argument("--input", "-i", required=True, help="输入 Markdown 文件路径")
    parser.add_argument("--output", "-o", required=True, help="输出 PPTX 文件路径")
    parser.add_argument("--title", "-t", default="", help="PPT 标题")
    parser.add_argument(
        "--theme",
        default="tech_blue",
        choices=list(THEMES.keys()),
        help="配色方案 (默认: tech_blue)",
    )
    args = parser.parse_args()
    generate_ppt(args.input, args.output, args.title, args.theme)


if __name__ == "__main__":
    main()
